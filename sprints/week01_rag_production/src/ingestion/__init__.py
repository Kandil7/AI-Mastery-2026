"""
Production-Grade Document Ingestion Module for RAG Systems

This module implements robust document ingestion capabilities supporting multiple
file formats commonly found in enterprise environments. The ingestion pipeline
handles various document types with appropriate parsing, text extraction,
and metadata enrichment to ensure high-quality input for the RAG system.

Key Features:
- Multi-format support (PDF, DOCX, PPTX, XLSX, TXT, MD, HTML)
- Robust error handling and fallback mechanisms
- Metadata extraction and enrichment
- Content validation and cleaning
- Progress tracking and logging
- Secure file processing with size limits

Supported Formats:
- PDF: Using PyMuPDF for high-quality text extraction
- DOCX: Using python-docx for Word documents
- PPTX: Using python-pptx for PowerPoint presentations
- XLSX: Using pandas for Excel spreadsheets
- TXT: Plain text files
- MD: Markdown files
- HTML: Web pages and HTML documents

Security Considerations:
- File size limits to prevent resource exhaustion
- Safe parsing to prevent malicious content execution
- Content validation to ensure quality input

Example:
    >>> from src.ingestion import IngestionPipeline
    >>> 
    >>> pipeline = IngestionPipeline(max_file_size=10*1024*1024)  # 10MB limit
    >>> documents = pipeline.load_from_directory("./documents/")
    >>> print(f"Loaded {len(documents)} documents")
"""

import os
import hashlib
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
from datetime import datetime
import re

# Import optional dependencies with fallbacks
try:
    import fitz  # PyMuPDF for PDF processing
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False
    print("PyMuPDF not available. PDF processing will be limited.")

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("python-docx not available. DOCX processing will be disabled.")

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False
    print("pandas not available. XLSX processing will be disabled.")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("python-pptx not available. PPTX processing will be disabled.")

# Import core dependencies
from src.retrieval import Document


class IngestionPipeline:
    """
    Main ingestion pipeline for loading documents from various sources.
    
    This class orchestrates the document loading process, handling different
    file formats and applying appropriate parsing strategies. It includes
    validation, metadata extraction, and error handling to ensure robust
    ingestion in production environments.
    
    Args:
        max_file_size (int): Maximum allowed file size in bytes (default 10MB)
        chunk_size (int): Size of text chunks for processing (default 1000 chars)
        enable_ocr (bool): Whether to enable OCR for image-based PDFs (requires additional setup)
        
    Example:
        >>> pipeline = IngestionPipeline(max_file_size=5*1024*1024)  # 5MB limit
        >>> docs = pipeline.load_from_directory("./docs/", recursive=True)
        >>> print(f"Loaded {len(docs)} documents")
    """
    
    def __init__(self, max_file_size: int = 10 * 1024 * 1024, chunk_size: int = 1000, enable_ocr: bool = False):
        """
        Initialize the ingestion pipeline with configuration options.
        
        Args:
            max_file_size (int): Maximum allowed file size in bytes
            chunk_size (int): Size of text chunks for processing
            enable_ocr (bool): Whether to enable OCR for image-based PDFs
        """
        self.max_file_size = max_file_size
        self.chunk_size = chunk_size
        self.enable_ocr = enable_ocr
        self.logger = logging.getLogger(__name__)
        
        # Supported file extensions
        self.supported_extensions = {
            '.pdf': self._load_pdf,
            '.txt': self._load_txt,
            '.md': self._load_md,
            '.html': self._load_html,
            '.htm': self._load_html,
        }
        
        # Add optional formats if dependencies are available
        if HAS_DOCX:
            self.supported_extensions.update({'.docx': self._load_docx})
        if HAS_PPTX:
            self.supported_extensions.update({'.pptx': self._load_pptx})
        if HAS_PANDAS:
            self.supported_extensions.update({'.xlsx': self._load_xlsx, '.xls': self._load_xlsx})
    
    def load_from_file(self, file_path: Union[str, Path]) -> List[Document]:
        """
        Load documents from a single file.
        
        Args:
            file_path (Union[str, Path]): Path to the file to load
            
        Returns:
            List[Document]: List of Document objects extracted from the file
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File does not exist: {file_path}")
        
        if file_path.stat().st_size > self.max_file_size:
            raise ValueError(f"File size exceeds limit: {file_path} ({file_path.stat().st_size} bytes)")
        
        ext = file_path.suffix.lower()
        if ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file format: {ext}. Supported: {list(self.supported_extensions.keys())}")
        
        try:
            loader_func = self.supported_extensions[ext]
            return loader_func(file_path)
        except Exception as e:
            self.logger.error(f"Error loading file {file_path}: {str(e)}")
            raise
    
    def load_from_directory(self, directory: Union[str, Path], recursive: bool = True) -> List[Document]:
        """
        Load documents from all supported files in a directory.
        
        Args:
            directory (Union[str, Path]): Directory to scan for documents
            recursive (bool): Whether to scan subdirectories recursively
            
        Returns:
            List[Document]: List of Document objects from all files
        """
        directory = Path(directory)
        if not directory.is_dir():
            raise ValueError(f"Directory does not exist: {directory}")
        
        documents = []
        file_pattern = "**/*" if recursive else "*"
        
        for ext in self.supported_extensions.keys():
            for file_path in directory.glob(file_pattern + ext):
                try:
                    docs = self.load_from_file(file_path)
                    documents.extend(docs)
                    self.logger.info(f"Loaded {len(docs)} documents from {file_path}")
                except Exception as e:
                    self.logger.warning(f"Failed to load {file_path}: {str(e)}")
                    continue
        
        return documents
    
    def _load_pdf(self, file_path: Path) -> List[Document]:
        """
        Load documents from a PDF file.
        
        Args:
            file_path (Path): Path to the PDF file
            
        Returns:
            List[Document]: List of Document objects, one per page
        """
        if not HAS_FITZ:
            raise ImportError("PyMuPDF is required to process PDF files")
        
        documents = []
        doc = fitz.open(file_path)
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            
            # Clean up text
            text = self._clean_text(text)
            
            if text.strip():  # Only add if text is not empty
                doc_id = f"{file_path.stem}_page_{page_num + 1}_{hashlib.md5(text.encode()).hexdigest()[:8]}"
                
                document = Document(
                    id=doc_id,
                    content=text,
                    source=str(file_path),
                    doc_type="pdf",
                    created_at=datetime.fromtimestamp(file_path.stat().st_ctime).isoformat(),
                    updated_at=datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
                    page_number=page_num + 1,
                    metadata={
                        "file_path": str(file_path),
                        "file_size": file_path.stat().st_size,
                        "page_number": page_num + 1,
                        "total_pages": len(doc)
                    }
                )
                documents.append(document)
        
        doc.close()
        return documents
    
    def _load_txt(self, file_path: Path) -> List[Document]:
        """
        Load documents from a plain text file.
        
        Args:
            file_path (Path): Path to the text file
            
        Returns:
            List[Document]: List of Document objects
        """
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Clean up text
        content = self._clean_text(content)
        
        doc_id = f"{file_path.stem}_{hashlib.md5(content.encode()).hexdigest()[:8]}"
        
        document = Document(
            id=doc_id,
            content=content,
            source=str(file_path),
            doc_type="txt",
            created_at=datetime.fromtimestamp(file_path.stat().st_ctime).isoformat(),
            updated_at=datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
            metadata={
                "file_path": str(file_path),
                "file_size": file_path.stat().st_size
            }
        )
        
        return [document]
    
    def _load_md(self, file_path: Path) -> List[Document]:
        """
        Load documents from a Markdown file.
        
        Args:
            file_path (Path): Path to the Markdown file
            
        Returns:
            List[Document]: List of Document objects
        """
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Clean up text
        content = self._clean_text(content)
        
        doc_id = f"{file_path.stem}_{hashlib.md5(content.encode()).hexdigest()[:8]}"
        
        document = Document(
            id=doc_id,
            content=content,
            source=str(file_path),
            doc_type="markdown",
            created_at=datetime.fromtimestamp(file_path.stat().st_ctime).isoformat(),
            updated_at=datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
            metadata={
                "file_path": str(file_path),
                "file_size": file_path.stat().st_size
            }
        )
        
        return [document]
    
    def _load_html(self, file_path: Path) -> List[Document]:
        """
        Load documents from an HTML file.
        
        Args:
            file_path (Path): Path to the HTML file
            
        Returns:
            List[Document]: List of Document objects
        """
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("beautifulsoup4 is required to process HTML files")
        
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Parse HTML and extract text
        soup = BeautifulSoup(content, 'html.parser')
        text = soup.get_text()
        
        # Clean up text
        text = self._clean_text(text)
        
        doc_id = f"{file_path.stem}_{hashlib.md5(text.encode()).hexdigest()[:8]}"
        
        document = Document(
            id=doc_id,
            content=text,
            source=str(file_path),
            doc_type="html",
            created_at=datetime.fromtimestamp(file_path.stat().st_ctime).isoformat(),
            updated_at=datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
            metadata={
                "file_path": str(file_path),
                "file_size": file_path.stat().st_size
            }
        )
        
        return [document]
    
    def _load_docx(self, file_path: Path) -> List[Document]:
        """
        Load documents from a DOCX file.
        
        Args:
            file_path (Path): Path to the DOCX file
            
        Returns:
            List[Document]: List of Document objects
        """
        if not HAS_DOCX:
            raise ImportError("python-docx is required to process DOCX files")
        
        doc = DocxDocument(file_path)
        full_text = []
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            full_text.append(paragraph.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    full_text.append(cell.text)
        
        content = '\n'.join(full_text)
        
        # Clean up text
        content = self._clean_text(content)
        
        doc_id = f"{file_path.stem}_{hashlib.md5(content.encode()).hexdigest()[:8]}"
        
        document = Document(
            id=doc_id,
            content=content,
            source=str(file_path),
            doc_type="docx",
            created_at=datetime.fromtimestamp(file_path.stat().st_ctime).isoformat(),
            updated_at=datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
            metadata={
                "file_path": str(file_path),
                "file_size": file_path.stat().st_size
            }
        )
        
        return [document]
    
    def _load_pptx(self, file_path: Path) -> List[Document]:
        """
        Load documents from a PPTX file.
        
        Args:
            file_path (Path): Path to the PPTX file
            
        Returns:
            List[Document]: List of Document objects
        """
        if not HAS_PPTX:
            raise ImportError("python-pptx is required to process PPTX files")
        
        presentation = Presentation(file_path)
        slides_text = []
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            slide_content = '\n'.join(slide_text)
            if slide_content.strip():
                doc_id = f"{file_path.stem}_slide_{slide_num + 1}_{hashlib.md5(slide_content.encode()).hexdigest()[:8]}"
                
                document = Document(
                    id=doc_id,
                    content=slide_content,
                    source=str(file_path),
                    doc_type="pptx",
                    created_at=datetime.fromtimestamp(file_path.stat().st_ctime).isoformat(),
                    updated_at=datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
                    page_number=slide_num + 1,  # Using page_number for slide number
                    metadata={
                        "file_path": str(file_path),
                        "file_size": file_path.stat().st_size,
                        "slide_number": slide_num + 1,
                        "total_slides": len(presentation.slides)
                    }
                )
                slides_text.append(document)
        
        return slides_text
    
    def _load_xlsx(self, file_path: Path) -> List[Document]:
        """
        Load documents from an XLSX file.
        
        Args:
            file_path (Path): Path to the XLSX file
            
        Returns:
            List[Document]: List of Document objects
        """
        if not HAS_PANDAS:
            raise ImportError("pandas is required to process XLSX files")
        
        # Read all sheets in the Excel file
        excel_file = pd.ExcelFile(file_path)
        documents = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # Convert DataFrame to text representation
            content = f"Sheet: {sheet_name}\n\n{df.to_string()}"
            
            # Clean up text
            content = self._clean_text(content)
            
            if content.strip():
                doc_id = f"{file_path.stem}_{sheet_name}_{hashlib.md5(content.encode()).hexdigest()[:8]}"
                
                document = Document(
                    id=doc_id,
                    content=content,
                    source=str(file_path),
                    doc_type="xlsx",
                    created_at=datetime.fromtimestamp(file_path.stat().st_ctime).isoformat(),
                    updated_at=datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
                    section_title=sheet_name,
                    metadata={
                        "file_path": str(file_path),
                        "file_size": file_path.stat().st_size,
                        "sheet_name": sheet_name,
                        "rows": len(df),
                        "columns": len(df.columns)
                    }
                )
                documents.append(document)
        
        return documents
    
    def _clean_text(self, text: str) -> str:
        """
        Clean and normalize text content.
        
        Args:
            text (str): Raw text content
            
        Returns:
            str: Cleaned text content
        """
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters that might cause issues
        text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
        
        return text.strip()


class ChunkingStrategy:
    """
    Abstract base class for different chunking strategies.
    
    This class defines the interface for various chunking approaches
    that can be used to split documents into smaller, more manageable pieces.
    """
    
    def chunk(self, document: Document, chunk_size: int, overlap: int = 0) -> List[Document]:
        """
        Split a document into chunks according to the strategy.
        
        Args:
            document (Document): The document to chunk
            chunk_size (int): Target size of each chunk
            overlap (int): Number of characters to overlap between chunks
            
        Returns:
            List[Document]: List of chunked documents
        """
        raise NotImplementedError("Subclasses must implement the chunk method")


class FixedSizeChunker(ChunkingStrategy):
    """
    Fixed-size chunking strategy that splits documents at regular intervals.
    
    This approach divides documents into fixed-length segments regardless
    of semantic boundaries. It's simple but may split in the middle of
    meaningful text units.
    """
    
    def chunk(self, document: Document, chunk_size: int, overlap: int = 0) -> List[Document]:
        """
        Split a document into fixed-size chunks.
        
        Args:
            document (Document): The document to chunk
            chunk_size (int): Target size of each chunk
            overlap (int): Number of characters to overlap between chunks
            
        Returns:
            List[Document]: List of chunked documents
        """
        content = document.content
        chunks = []
        
        start = 0
        while start < len(content):
            end = start + chunk_size
            chunk_content = content[start:end]
            
            # Create a new document for this chunk
            chunk_doc = Document(
                id=f"{document.id}_chunk_{len(chunks) + 1}",
                content=chunk_content,
                source=document.source,
                doc_type=f"{document.doc_type}_chunk",
                created_at=document.created_at,
                updated_at=document.updated_at,
                access_control=document.access_control,
                page_number=document.page_number,
                section_title=document.section_title,
                metadata={**document.metadata, "chunk_index": len(chunks), "original_id": document.id}
            )
            
            chunks.append(chunk_doc)
            
            # Move to the next chunk position with overlap
            start = end - overlap if overlap < chunk_size else start + chunk_size
        
        return chunks


class SemanticChunker(ChunkingStrategy):
    """
    Semantic chunking strategy that respects document structure.
    
    This approach attempts to split documents at natural boundaries
    like paragraph breaks, headings, or sentences to maintain context
    and meaning within each chunk.
    """
    
    def chunk(self, document: Document, chunk_size: int, overlap: int = 0) -> List[Document]:
        """
        Split a document into semantically coherent chunks.
        
        Args:
            document (Document): The document to chunk
            chunk_size (int): Target size of each chunk
            overlap (int): Number of characters to overlap between chunks
            
        Returns:
            List[Document]: List of chunked documents
        """
        # Split by paragraphs first
        paragraphs = [p.strip() for p in document.content.split('\n') if p.strip()]
        
        chunks = []
        current_chunk = ""
        chunk_idx = 0
        
        for para in paragraphs:
            # If adding this paragraph would exceed chunk size
            if len(current_chunk) + len(para) > chunk_size and current_chunk:
                # Save the current chunk
                chunk_doc = Document(
                    id=f"{document.id}_chunk_{chunk_idx + 1}",
                    content=current_chunk.strip(),
                    source=document.source,
                    doc_type=f"{document.doc_type}_chunk",
                    created_at=document.created_at,
                    updated_at=document.updated_at,
                    access_control=document.access_control,
                    page_number=document.page_number,
                    section_title=document.section_title,
                    metadata={**document.metadata, "chunk_index": chunk_idx, "original_id": document.id}
                )
                chunks.append(chunk_doc)
                
                # Start a new chunk with potential overlap
                if overlap > 0:
                    # Take the end portion of the previous chunk for overlap
                    overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                    current_chunk = overlap_text + " " + para
                else:
                    current_chunk = para
                chunk_idx += 1
            else:
                # Add paragraph to current chunk
                if current_chunk:
                    current_chunk += "\n" + para
                else:
                    current_chunk = para
        
        # Add the last chunk if it has content
        if current_chunk.strip():
            chunk_doc = Document(
                id=f"{document.id}_chunk_{chunk_idx + 1}",
                content=current_chunk.strip(),
                source=document.source,
                doc_type=f"{document.doc_type}_chunk",
                created_at=document.created_at,
                updated_at=document.updated_at,
                access_control=document.access_control,
                page_number=document.page_number,
                section_title=document.section_title,
                metadata={**document.metadata, "chunk_index": chunk_idx, "original_id": document.id}
            )
            chunks.append(chunk_doc)
        
        return chunks


def create_ingestion_pipeline(max_file_size: int = 10 * 1024 * 1024, chunk_size: int = 1000) -> IngestionPipeline:
    """
    Factory function to create an ingestion pipeline with default configuration.
    
    Args:
        max_file_size (int): Maximum allowed file size in bytes
        chunk_size (int): Size of text chunks for processing
        
    Returns:
        IngestionPipeline: Configured ingestion pipeline
    """
    return IngestionPipeline(max_file_size=max_file_size, chunk_size=chunk_size)